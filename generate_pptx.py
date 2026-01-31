"""
JukeBoB AI - Professional PowerPoint Presentation Generator
Creates a polished executive presentation for Board of Directors
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Brand Colors
BRAND_YELLOW = RgbColor(245, 213, 71)
BRAND_GOLD = RgbColor(255, 215, 0)
DARK_BG = RgbColor(26, 26, 46)
DARK_BLUE = RgbColor(20, 30, 48)
WHITE = RgbColor(255, 255, 255)
LIGHT_GRAY = RgbColor(200, 200, 200)
SUCCESS_GREEN = RgbColor(72, 187, 120)
DANGER_RED = RgbColor(245, 101, 101)
JUKEBOX_PINK = RgbColor(240, 147, 251)
GAMES_CORAL = RgbColor(252, 182, 159)
AIDJ_CYAN = RgbColor(0, 212, 255)

def set_slide_background(slide, color):
    """Set solid background color for a slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_background(slide, DARK_BG)
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🎵 JUKEBOB"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = BRAND_YELLOW
    p.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.8))
    tf = tagline_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Your Ultimate Party Entertainment Hub"
    p.font.size = Pt(28)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(4.8), Inches(3), Inches(0.6))
    badge.fill.solid()
    badge.fill.fore_color.rgb = BRAND_YELLOW
    badge.line.fill.background()
    badge_tf = badge.text_frame
    badge_tf.paragraphs[0].text = "Executive Presentation"
    badge_tf.paragraphs[0].font.size = Pt(16)
    badge_tf.paragraphs[0].font.bold = True
    badge_tf.paragraphs[0].font.color.rgb = DARK_BG
    badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Board of Directors Briefing • December 2025"
    p.font.size = Pt(14)
    p.font.color.rgb = RgbColor(128, 128, 128)
    p.alignment = PP_ALIGN.CENTER

def add_what_is_jukebob(prs):
    """Slide 2: What is JukeBoB?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    
    # Section badge
    add_section_badge(slide, "INTRODUCTION", 0.5, 0.4)
    
    # Title
    add_slide_title(slide, "What is JukeBoB?", 0.5, 0.7)
    
    # Description
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.8))
    tf = desc_box.text_frame
    p = tf.paragraphs[0]
    p.text = "A comprehensive party entertainment platform combining music, games, and AI DJ into one unified experience."
    p.font.size = Pt(18)
    p.font.color.rgb = LIGHT_GRAY
    
    # Three pillars
    pillars = [
        ("🎵", "Jukebox", "Crowd-controlled music\nwith VIP tipping", JUKEBOX_PINK),
        ("🎮", "Games", "Real-time multiplayer\nwith emoji reactions", GAMES_CORAL),
        ("🎧", "AI DJ", "Intelligent dual-turntable\nmixing system", AIDJ_CYAN),
    ]
    
    x_start = 0.5
    for i, (emoji, title, desc, color) in enumerate(pillars):
        x = x_start + i * 3.2
        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.5), Inches(2.9), Inches(2.8))
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.fill.background()
        
        # Emoji
        emoji_box = slide.shapes.add_textbox(Inches(x), Inches(2.7), Inches(2.9), Inches(0.8))
        tf = emoji_box.text_frame
        p = tf.paragraphs[0]
        p.text = emoji
        p.font.size = Pt(48)
        p.alignment = PP_ALIGN.CENTER
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(x), Inches(3.5), Inches(2.9), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE if color != GAMES_CORAL else DARK_BG
        p.alignment = PP_ALIGN.CENTER
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(4.1), Inches(2.7), Inches(1))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE if color != GAMES_CORAL else DARK_BG
        p.alignment = PP_ALIGN.CENTER

def add_platform_metrics(prs):
    """Slide 3: Platform at a Glance"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    
    add_section_badge(slide, "PLATFORM OVERVIEW", 0.5, 0.4)
    add_slide_title(slide, "Platform at a Glance", 0.5, 0.7)
    
    metrics = [
        ("5%", "Platform Fee"),
        ("3", "Entertainment Modules"),
        ("₹10+", "VIP Threshold"),
        ("3", "Platforms"),
        ("Real-time", "WebSocket Sync"),
        ("5,200+", "Lines of Code"),
    ]
    
    for i, (value, label) in enumerate(metrics):
        row = i // 3
        col = i % 3
        x = 0.5 + col * 3.2
        y = 1.8 + row * 2.0
        
        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.9), Inches(1.6))
        card.fill.solid()
        card.fill.fore_color.rgb = RgbColor(40, 40, 70)
        card.line.color.rgb = BRAND_YELLOW
        
        # Value
        val_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.2), Inches(2.9), Inches(0.8))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = BRAND_YELLOW
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        lbl_box = slide.shapes.add_textbox(Inches(x), Inches(y + 1.0), Inches(2.9), Inches(0.4))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER

def add_problem_solution(prs):
    """Slide 4: Problem & Solution"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    
    add_section_badge(slide, "MARKET OPPORTUNITY", 0.5, 0.4)
    add_slide_title(slide, "The Problem We Solve", 0.5, 0.7)
    
    # Left column - Problems
    prob_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(0.5))
    tf = prob_title.text_frame
    p = tf.paragraphs[0]
    p.text = "❌ Pain Points"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DANGER_RED
    
    problems = [
        "🎵  DJs decide all music without guest input",
        "😴  Boring waiting periods at events",
        "💸  Hosts miss tipping opportunities",
        "📱  Fragmented tools for entertainment",
    ]
    
    for i, prob in enumerate(problems):
        box = slide.shapes.add_textbox(Inches(0.5), Inches(2.1 + i * 0.6), Inches(4.5), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = prob
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
    
    # Right column - Solutions
    sol_title = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.5), Inches(0.5))
    tf = sol_title.text_frame
    p = tf.paragraphs[0]
    p.text = "✓ Our Solution"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_GREEN
    
    solutions = [
        "🗳️  Democratic music selection via tipping",
        "🎮  Games keep guests entertained",
        "💰  Passive income stream for hosts",
        "🎧  AI-powered DJ for seamless mixing",
    ]
    
    for i, sol in enumerate(solutions):
        box = slide.shapes.add_textbox(Inches(5.2), Inches(2.1 + i * 0.6), Inches(4.5), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = sol
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE

def add_jukebox_slide(prs):
    """Slide 5: Jukebox Module"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    
    add_section_badge(slide, "CORE MODULES", 0.5, 0.4)
    add_slide_title(slide, "🎵 Jukebox — Revenue Engine", 0.5, 0.7)
    
    # Left column
    left_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(0.4))
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = "VIP Queue System"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    features = [
        "👑  Tips ≥₹10 = VIP Priority",
        "📊  Sorted by tip amount (highest first)",
        "🎵  Regular queue: First-come-first-served",
        "🔄  Real-time sync for all users",
    ]
    
    for i, feat in enumerate(features):
        box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0 + i * 0.5), Inches(4.5), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = feat
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
    
    # Host features
    host_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(4.5), Inches(0.4))
    tf = host_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Host Features"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    host_feats = ["📱  QR code for easy joining", "💵  Dashboard with earnings", "✅  Mark songs played/skipped"]
    for i, feat in enumerate(host_feats):
        box = slide.shapes.add_textbox(Inches(0.5), Inches(4.7 + i * 0.45), Inches(4.5), Inches(0.45))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = feat
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
    
    # Right - Code box
    code_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.5), Inches(4.3), Inches(2.5))
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RgbColor(30, 30, 63)
    code_box.line.color.rgb = BRAND_YELLOW
    
    code_text = slide.shapes.add_textbox(Inches(5.4), Inches(1.7), Inches(4), Inches(2.2))
    tf = code_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "# Queue Priority Logic\n\nvip_songs = sorted(\n    [r for r in queue if tip >= 10],\n    key=lambda x: x.tip,\n    reverse=True\n)\n\nfinal_queue = vip + regular"
    p.font.size = Pt(12)
    p.font.color.rgb = BRAND_YELLOW
    p.font.name = "Consolas"

def add_games_slide(prs):
    """Slide 6: Games Module"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    
    add_section_badge(slide, "CORE MODULES", 0.5, 0.4)
    add_slide_title(slide, "🎮 Games — Engagement Engine", 0.5, 0.7)
    
    # Left column
    left_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(0.4))
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Tic Tac Toe Features"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    features = [
        "⚡  Real-time WebSocket gameplay",
        "🏆  Persistent scoreboard & leaderboard",
        "🔥  Win streak tracking",
        "🔄  Turn alternation on restart",
    ]
    
    for i, feat in enumerate(features):
        box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0 + i * 0.5), Inches(4.5), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = feat
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
    
    # Emoji reactions
    emoji_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(4.5), Inches(0.4))
    tf = emoji_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Emoji Reactions"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    emoji_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.7), Inches(4.5), Inches(0.8))
    tf = emoji_box.text_frame
    p = tf.paragraphs[0]
    p.text = "😂 😭 😎 😉 😘 😜 😱 👏 🌹 🏆"
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    
    # Right - Stats table
    stats_title = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.5), Inches(0.4))
    tf = stats_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Statistics Tracked"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    stats = [("Wins", "Total victories"), ("Losses", "Total defeats"), ("Draws", "Total ties"), ("Streak", "Current win streak"), ("Max Streak", "Best ever")]
    
    for i, (stat, desc) in enumerate(stats):
        row = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.2), Inches(2.0 + i * 0.55), Inches(4.3), Inches(0.5))
        row.fill.solid()
        row.fill.fore_color.rgb = RgbColor(40, 40, 70) if i % 2 == 0 else RgbColor(50, 50, 80)
        row.line.fill.background()
        
        stat_box = slide.shapes.add_textbox(Inches(5.4), Inches(2.05 + i * 0.55), Inches(1.5), Inches(0.4))
        tf = stat_box.text_frame
        p = tf.paragraphs[0]
        p.text = stat
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = BRAND_YELLOW
        
        desc_box = slide.shapes.add_textbox(Inches(7), Inches(2.05 + i * 0.55), Inches(2.5), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_GRAY

def add_aidj_slide(prs):
    """Slide 7: AI DJ Module"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    
    add_section_badge(slide, "CORE MODULES", 0.5, 0.4)
    add_slide_title(slide, "🎧 AI DJ — Innovation Engine", 0.5, 0.7)
    
    # Left column
    left_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(0.4))
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Dual Turntable System"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    features = ["🎛️  Deck A & Deck B audio players", "📤  Upload audio files directly", "🔊  Independent volume controls", "🔀  AI-powered crossfader"]
    
    for i, feat in enumerate(features):
        box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0 + i * 0.5), Inches(4.5), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = feat
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
    
    ai_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(4.5), Inches(0.4))
    tf = ai_title.text_frame
    p = tf.paragraphs[0]
    p.text = "AI Capabilities"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    ai_feats = ["🎭  Crowd mood analysis", "⏱️  3-second crossfade transitions", "💡  Smart song suggestions"]
    for i, feat in enumerate(ai_feats):
        box = slide.shapes.add_textbox(Inches(0.5), Inches(4.7 + i * 0.45), Inches(4.5), Inches(0.45))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = feat
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
    
    # Right - Mood detection
    mood_title = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.5), Inches(0.4))
    tf = mood_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Mood Detection"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    code_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(2.0), Inches(4.3), Inches(2.2))
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RgbColor(30, 30, 63)
    code_box.line.color.rgb = AIDJ_CYAN
    
    code_text = slide.shapes.add_textbox(Inches(5.4), Inches(2.2), Inches(4), Inches(2))
    tf = code_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "if avg_votes > 10:\n    return \"hyped\" 🔥\nelif avg_votes > 5:\n    return \"energetic\" ⚡\nelse:\n    return \"chill\" 😌"
    p.font.size = Pt(14)
    p.font.color.rgb = AIDJ_CYAN
    p.font.name = "Consolas"

def add_revenue_slide(prs):
    """Slide 8: Revenue Model"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    
    add_section_badge(slide, "MONETIZATION", 0.5, 0.4)
    add_slide_title(slide, "💰 Revenue Model", 0.5, 0.7)
    
    # Revenue flow boxes
    boxes_data = [
        (1, "₹100", "Guest Tips", RgbColor(60, 60, 90), WHITE),
        (3.8, "₹5", "Platform (5%)", BRAND_YELLOW, DARK_BG),
        (6.6, "₹95", "Host Gets", RgbColor(60, 60, 90), WHITE),
    ]
    
    for x, amount, label, bg_color, text_color in boxes_data:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.4), Inches(2.5), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = bg_color
        box.line.fill.background()
        
        amt_box = slide.shapes.add_textbox(Inches(x), Inches(1.5), Inches(2.5), Inches(0.6))
        tf = amt_box.text_frame
        p = tf.paragraphs[0]
        p.text = amount
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER
        
        lbl_box = slide.shapes.add_textbox(Inches(x), Inches(2.1), Inches(2.5), Inches(0.4))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER
    
    # Arrows
    for x in [3.5, 6.3]:
        arrow = slide.shapes.add_textbox(Inches(x), Inches(1.6), Inches(0.4), Inches(0.6))
        tf = arrow.text_frame
        p = tf.paragraphs[0]
        p.text = "→"
        p.font.size = Pt(32)
        p.font.color.rgb = BRAND_YELLOW
        p.alignment = PP_ALIGN.CENTER
    
    # Consumer Protection Section (LEFT)
    protection_title = slide.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(4.5), Inches(0.4))
    tf = protection_title.text_frame
    p = tf.paragraphs[0]
    p.text = "🛡️ Consumer Protection"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    protection_items = [
        "💰  Tips held in escrow until song plays",
        "🔄  Full refund if DJ skips the song",
        "❌  DJ gets ₹0 for unplayed requests",
        "✅  Automatic refunds, no support needed",
    ]
    
    for i, item in enumerate(protection_items):
        box = slide.shapes.add_textbox(Inches(0.5), Inches(3.4 + i * 0.45), Inches(4.5), Inches(0.45))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
    
    # Future Revenue Section (RIGHT)
    table_title = slide.shapes.add_textbox(Inches(5.2), Inches(2.9), Inches(4.5), Inches(0.4))
    tf = table_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Future Revenue Streams"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    streams = [
        ("Transaction Fees (5%)", "✅ Live"),
        ("Premium Features", "🔜 Planned"),
        ("White-label Licensing", "🔜 Planned"),
        ("Subscription Model", "🔜 Planned"),
    ]
    
    for i, (stream, status) in enumerate(streams):
        y = 3.4 + i * 0.45
        row_color = RgbColor(40, 40, 70) if i % 2 == 0 else RgbColor(50, 50, 80)
        
        row = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.2), Inches(y), Inches(4.3), Inches(0.4))
        row.fill.solid()
        row.fill.fore_color.rgb = row_color
        row.line.fill.background()
        
        stream_txt = slide.shapes.add_textbox(Inches(5.3), Inches(y + 0.05), Inches(2.8), Inches(0.35))
        tf = stream_txt.text_frame
        p = tf.paragraphs[0]
        p.text = stream
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        
        status_txt = slide.shapes.add_textbox(Inches(8), Inches(y + 0.05), Inches(1.5), Inches(0.35))
        tf = status_txt.text_frame
        p = tf.paragraphs[0]
        p.text = status
        p.font.size = Pt(11)
        p.font.color.rgb = SUCCESS_GREEN if "✅" in status else RgbColor(255, 165, 0)


def add_tech_architecture(prs):
    """Slide 9: Technical Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    
    add_section_badge(slide, "TECHNOLOGY", 0.5, 0.4)
    add_slide_title(slide, "🏗️ Technical Architecture", 0.5, 0.7)
    
    tech_boxes = [
        ("⚡ Backend", ["Python 3.11", "FastAPI Framework", "WebSocket Server", "821 lines of code"], 0.5, 1.5),
        ("🌐 Web Frontend", ["HTML5 + CSS3", "Vanilla JavaScript", "Multi-theme Design", "2,500+ lines"], 5, 1.5),
        ("📱 Mobile App", ["Flutter Framework", "Dart Language", "iOS + Android", "Material 3 Design"], 0.5, 3.5),
        ("🔒 Security", ["SHA256 Hashing", "Session Auth", "CORS Protection", "Input Validation"], 5, 3.5),
    ]
    
    for title, items, x, y in tech_boxes:
        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.3), Inches(1.7))
        box.fill.solid()
        box.fill.fore_color.rgb = RgbColor(35, 35, 60)
        box.line.color.rgb = BRAND_YELLOW
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.1), Inches(4), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BRAND_YELLOW
        
        # Items
        for i, item in enumerate(items):
            item_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.5 + i * 0.3), Inches(4), Inches(0.3))
            tf = item_box.text_frame
            p = tf.paragraphs[0]
            p.text = "• " + item
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE

def add_api_overview(prs):
    """Slide 10: API Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    
    add_section_badge(slide, "TECHNOLOGY", 0.5, 0.4)
    add_slide_title(slide, "🔌 API Architecture", 0.5, 0.7)
    
    # Metrics
    metrics = [("30+", "REST Endpoints"), ("1", "WebSocket Channel"), ("50+", "Functions")]
    for i, (val, lbl) in enumerate(metrics):
        x = 1 + i * 3
        
        val_box = slide.shapes.add_textbox(Inches(x), Inches(1.5), Inches(2.5), Inches(0.8))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = BRAND_YELLOW
        p.alignment = PP_ALIGN.CENTER
        
        lbl_box = slide.shapes.add_textbox(Inches(x), Inches(2.3), Inches(2.5), Inches(0.4))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = lbl
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER
    
    # Table
    endpoints = [
        ("Sessions", "4", "Create, Join, Resume, Get"),
        ("Requests", "5", "Submit, Complete, Skip, Tip, List"),
        ("Games", "7", "Create, Join, Move, Emoji, Restart"),
        ("Payments", "3", "Checkout, Transaction, Revenue"),
        ("AI DJ", "3", "Enable, Playlist, Next Track"),
    ]
    
    headers = ["Module", "Endpoints", "Key Functions"]
    for j, h in enumerate(headers):
        w = 2.5 if j < 2 else 4.5
        hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5 + sum([2.5, 2.5][:j]) if j < 2 else 5), Inches(3), Inches(w), Inches(0.45))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = RgbColor(60, 50, 30)
        hdr.line.fill.background()
        
        txt = slide.shapes.add_textbox(Inches(0.6 + sum([2.5, 2.5][:j]) if j < 2 else 5.1), Inches(3.05), Inches(w - 0.2), Inches(0.4))
        tf = txt.text_frame
        p = tf.paragraphs[0]
        p.text = h
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = BRAND_YELLOW
    
    for i, (mod, cnt, funcs) in enumerate(endpoints):
        y = 3.5 + i * 0.5
        row_color = RgbColor(40, 40, 70) if i % 2 == 0 else RgbColor(50, 50, 80)
        
        for j, (val, w) in enumerate([(mod, 2.5), (cnt, 2.5), (funcs, 4.5)]):
            x_offset = 0.5 + sum([2.5, 2.5][:j]) if j < 2 else 5
            cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_offset), Inches(y), Inches(w), Inches(0.45))
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_color
            cell.line.fill.background()
            
            txt = slide.shapes.add_textbox(Inches(x_offset + 0.1), Inches(y + 0.05), Inches(w - 0.2), Inches(0.4))
            tf = txt.text_frame
            p = tf.paragraphs[0]
            p.text = val
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE

def add_roadmap(prs):
    """Slide 11: Growth Roadmap"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    
    add_section_badge(slide, "FUTURE", 0.5, 0.4)
    add_slide_title(slide, "🚀 Growth Roadmap", 0.5, 0.7)
    
    phases = [
        ("Q1", "Production Ready", ["✓ Supabase database integration", "✓ Stripe + UPI payment gateway", "✓ User registration & authentication"]),
        ("Q2", "Feature Enhancement", ["○ More games (Mafia, Trivia)", "○ Advanced AI DJ with beat matching", "○ Push notifications"]),
        ("Q3-Q4", "Scale", ["○ Kubernetes deployment", "○ Analytics dashboard", "○ White-label solution"]),
    ]
    
    for i, (phase, title, items) in enumerate(phases):
        y = 1.5 + i * 1.7
        
        # Phase badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(1.2), Inches(0.5))
        badge.fill.solid()
        badge.fill.fore_color.rgb = BRAND_YELLOW
        badge.line.fill.background()
        
        badge_txt = slide.shapes.add_textbox(Inches(0.5), Inches(y + 0.05), Inches(1.2), Inches(0.4))
        tf = badge_txt.text_frame
        p = tf.paragraphs[0]
        p.text = phase
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = DARK_BG
        p.alignment = PP_ALIGN.CENTER
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1.9), Inches(y), Inches(4), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        # Items
        for j, item in enumerate(items):
            item_box = slide.shapes.add_textbox(Inches(1.9), Inches(y + 0.5 + j * 0.35), Inches(7), Inches(0.35))
            tf = item_box.text_frame
            p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(13)
            p.font.color.rgb = SUCCESS_GREEN if item.startswith("✓") else RgbColor(255, 165, 0)

def add_competitive_advantages(prs):
    """Slide 12: Competitive Advantages"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    
    add_section_badge(slide, "STRATEGY", 0.5, 0.4)
    add_slide_title(slide, "🏆 Competitive Advantages", 0.5, 0.7)
    
    advantages = [
        ("🥇", "First-Mover Advantage", "No direct competitors in integrated party entertainment"),
        ("🧩", "Three-in-One Platform", "Music, games, and DJ combined—reduces user friction"),
        ("⚡", "Real-Time Technology", "WebSocket creates sticky, engaging experiences"),
        ("💰", "Day-1 Revenue", "5% transaction fee built into platform from launch"),
    ]
    
    for i, (emoji, title, desc) in enumerate(advantages):
        row = i // 2
        col = i % 2
        x = 0.5 + col * 4.8
        y = 1.5 + row * 2.2
        
        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.5), Inches(1.8))
        card.fill.solid()
        card.fill.fore_color.rgb = RgbColor(40, 40, 70)
        card.line.color.rgb = BRAND_YELLOW
        
        # Emoji
        emoji_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.15), Inches(0.6), Inches(0.6))
        tf = emoji_box.text_frame
        p = tf.paragraphs[0]
        p.text = emoji
        p.font.size = Pt(28)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(x + 0.8), Inches(y + 0.2), Inches(3.5), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = BRAND_YELLOW
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.8), Inches(4.1), Inches(0.8))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_GRAY

def add_summary(prs):
    """Slide 13: Key Takeaways"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    
    add_section_badge(slide, "SUMMARY", 0.5, 0.4)
    add_slide_title(slide, "📋 Key Takeaways", 0.5, 0.7)
    
    # Highlight box
    highlight = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.4), Inches(7), Inches(1.1))
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = BRAND_YELLOW
    highlight.line.fill.background()
    
    h_title = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(7), Inches(0.5))
    tf = h_title.text_frame
    p = tf.paragraphs[0]
    p.text = "JukeBoB is Production-Ready"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_BG
    p.alignment = PP_ALIGN.CENTER
    
    h_sub = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(7), Inches(0.4))
    tf = h_sub.text_frame
    p = tf.paragraphs[0]
    p.text = "Fully functional MVP with clear path to scale"
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_BG
    p.alignment = PP_ALIGN.CENTER
    
    # Two columns
    # Left - What we have
    left_title = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(4.5), Inches(0.5))
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = "✅ What We Have"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_GREEN
    
    have_items = ["Web + iOS + Android applications", "Three entertainment modules", "Real-time synchronization", "Built-in 5% monetization", "VIP priority queue system"]
    for i, item in enumerate(have_items):
        box = slide.shapes.add_textbox(Inches(0.5), Inches(3.4 + i * 0.45), Inches(4.5), Inches(0.45))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = "✓  " + item
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
    
    # Right - What's next
    right_title = slide.shapes.add_textbox(Inches(5.2), Inches(2.8), Inches(4.5), Inches(0.5))
    tf = right_title.text_frame
    p = tf.paragraphs[0]
    p.text = "🔜 What's Next"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 165, 0)
    
    next_items = ["Payment gateway integration", "Database persistence", "User authentication", "More games & features", "Scale infrastructure"]
    for i, item in enumerate(next_items):
        box = slide.shapes.add_textbox(Inches(5.2), Inches(3.4 + i * 0.45), Inches(4.5), Inches(0.45))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = "○  " + item
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE

def add_thank_you(prs):
    """Slide 14: Thank You"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🎵 Thank You"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = BRAND_YELLOW
    p.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.6))
    tf = tagline_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions & Discussion"
    p.font.size = Pt(28)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Brand info
    brand_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(9), Inches(0.8))
    tf = brand_box.text_frame
    p = tf.paragraphs[0]
    p.text = "JukeBoB AI"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BRAND_YELLOW
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "Your Ultimate Party Entertainment Hub"
    p2.font.size = Pt(16)
    p2.font.color.rgb = RgbColor(128, 128, 128)
    p2.alignment = PP_ALIGN.CENTER
    
    # Date badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(5.3), Inches(3), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = BRAND_YELLOW
    badge.line.fill.background()
    
    badge_txt = slide.shapes.add_textbox(Inches(3.5), Inches(5.35), Inches(3), Inches(0.4))
    tf = badge_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "December 2025"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_BG
    p.alignment = PP_ALIGN.CENTER

def add_section_badge(slide, text, x, y):
    """Add a section badge"""
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2), Inches(0.35))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RgbColor(60, 50, 30)
    badge.line.fill.background()
    
    txt = slide.shapes.add_textbox(Inches(x), Inches(y + 0.02), Inches(2), Inches(0.3))
    tf = txt.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = BRAND_YELLOW
    p.alignment = PP_ALIGN.CENTER

def add_slide_title(slide, text, x, y):
    """Add slide title"""
    title_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BRAND_YELLOW

def main():
    """Generate the PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Add all slides
    add_title_slide(prs)
    add_what_is_jukebob(prs)
    add_platform_metrics(prs)
    add_problem_solution(prs)
    add_jukebox_slide(prs)
    add_games_slide(prs)
    add_aidj_slide(prs)
    add_revenue_slide(prs)
    add_tech_architecture(prs)
    add_api_overview(prs)
    add_roadmap(prs)
    add_competitive_advantages(prs)
    add_summary(prs)
    add_thank_you(prs)
    
    # Save
    output_path = "JukeBoB_Executive_Presentation.pptx"
    prs.save(output_path)
    print(f"[OK] Presentation saved: {output_path}")
    print(f"[INFO] Total slides: 14")

if __name__ == "__main__":
    main()
